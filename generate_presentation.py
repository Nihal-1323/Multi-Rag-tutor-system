"""
PowerPoint Presentation Generator for Multi-Modal RAG System
Run: pip install python-pptx
Then: python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Multi-Modal RAG System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(99, 102, 241)
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Intelligent Education Tutor with Hybrid Retrieval"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: System Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "An AI-powered education assistant that:"
    
    points = [
        "Processes multi-modal learning materials (text, PDF, images, audio)",
        "Uses hybrid retrieval combining vector search and knowledge graphs",
        "Provides accurate answers with source attribution",
        "Tracks performance metrics in real-time",
        "Offers two distinct user interfaces"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 3: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    components = [
        ("Frontend", "React + TypeScript, Two themes (Professional & Terminal)"),
        ("Backend", "FastAPI (Python), RESTful API"),
        ("Vector Database", "Weaviate - Semantic search with embeddings"),
        ("Graph Database", "Neo4j - Concept relationships and traversal"),
        ("LLM", "Ollama - Answer generation"),
        ("Deployment", "Docker Compose - Containerized services")
    ]
    
    for comp, desc in components:
        p = tf.add_paragraph()
        p.text = f"{comp}: {desc}"
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 4: Hybrid Retrieval Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Hybrid Retrieval Pipeline"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    steps = [
        ("1. Query Understanding", "Analyze user question and extract key concepts"),
        ("2. Dual Retrieval", "Vector search + Graph traversal in parallel"),
        ("3. Fusion", "Combine results from both methods"),
        ("4. Reranking", "Cross-encoder model improves relevance"),
        ("5. Answer Generation", "LLM creates response with citations")
    ]
    
    for step, desc in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(99, 102, 241)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.level = 1
        p2.font.size = Pt(16)
    
    # Slide 5: Document Processing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Document Processing Workflow"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    workflow = [
        "Upload → User uploads learning materials",
        "Extract → Content extraction from various formats",
        "Chunk → Split into manageable pieces",
        "Embed → Generate vector embeddings",
        "Store → Save to Weaviate vector database",
        "Graph → Build concept relationships in Neo4j"
    ]
    
    for item in workflow:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Slide 6: Ranking Algorithm
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Intelligent Ranking Algorithm"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Multi-factor scoring system:"
    
    factors = [
        "Exact Match: +100 points",
        "Word Frequency: +5 per occurrence",
        "Proximity Bonus: +20 for nearby terms",
        "Cross-Encoder Reranking: Deep semantic relevance",
        "Final Sorting: Highest score first"
    ]
    
    for factor in factors:
        p = tf.add_paragraph()
        p.text = factor
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 7: Performance Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Metrics"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    metrics = [
        ("Precision@K", "Accuracy of top-K retrieved documents"),
        ("Recall@K", "Coverage of relevant documents"),
        ("F1 Score", "Harmonic mean of precision and recall"),
        ("Latency", "Response time in milliseconds"),
        ("Vector Search", "Embedding similarity performance"),
        ("Graph Grounding", "Knowledge graph contribution")
    ]
    
    for metric, desc in metrics:
        p = tf.add_paragraph()
        p.text = f"{metric}: {desc}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # Slide 8: Two Frontend Interfaces
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Dual Frontend Experience"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Port 3000 - EduCore (Professional)"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 51, 234)
    
    features1 = ["Black theme with purple accents", "Sidebar navigation", "Card-based metrics", "Modern, clean design"]
    for f in features1:
        p = tf.add_paragraph()
        p.text = f
        p.level = 1
        p.font.size = Pt(16)
    
    tf.add_paragraph()
    
    p = tf.add_paragraph()
    p.text = "Port 3001 - Terminal (Developer)"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(20, 184, 166)
    
    features2 = ["Dark teal/cyan theme", "Horizontal tabs", "Terminal-style interface", "Technical aesthetic"]
    for f in features2:
        p = tf.add_paragraph()
        p.text = f
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 9: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    stack = [
        ("Backend", "FastAPI, Python, Uvicorn"),
        ("Databases", "Weaviate (Vector), Neo4j (Graph)"),
        ("AI/ML", "Ollama LLM, Multi-modal embeddings"),
        ("Frontend", "React, TypeScript, Vite, TailwindCSS"),
        ("Visualization", "Force-Graph, React-Markdown"),
        ("Deployment", "Docker, Docker Compose")
    ]
    
    for category, tech in stack:
        p = tf.add_paragraph()
        p.text = f"{category}: {tech}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    # Slide 10: Key Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Benefits"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    benefits = [
        "🎯 Accuracy: Hybrid approach improves retrieval quality",
        "🔍 Transparency: Shows ranking process and sources",
        "⚡ Performance: Real-time metrics tracking",
        "🎨 Flexibility: Multiple content types supported",
        "📊 Insights: Visual knowledge graph representation",
        "🚀 Scalability: Docker-based deployment",
        "💡 User-Friendly: Two distinct interface options"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Slide 11: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Use Cases"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    cases = [
        "📚 Education: Student learning assistant for course materials",
        "🏢 Enterprise: Internal knowledge base Q&A system",
        "🔬 Research: Literature review and paper analysis",
        "📖 Documentation: Technical documentation assistant",
        "🎓 Training: Corporate training material helper",
        "💼 Onboarding: New employee knowledge access"
    ]
    
    for case in cases:
        p = tf.add_paragraph()
        p.text = case
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Slide 12: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    enhancements = [
        "Multi-language support",
        "Advanced audio/video processing",
        "Collaborative learning features",
        "Personalized learning paths",
        "Integration with LMS platforms",
        "Mobile application",
        "Advanced analytics dashboard"
    ]
    
    for enhancement in enhancements:
        p = tf.add_paragraph()
        p.text = enhancement
        p.font.size = Pt(20)
        p.space_after = Pt(14)
    
    # Slide 13: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = RGBColor(99, 102, 241)
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Multi-Modal RAG System\nIntelligent Education Platform"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(20)
    contact_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('MultiModal_RAG_System_Presentation.pptx')
    print("✅ Presentation created: MultiModal_RAG_System_Presentation.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("Run: pip install python-pptx")
